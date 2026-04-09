from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = title_text
    subtitle.text = subtitle_text
    
    # Simple styling
    title.text_frame.paragraphs[0].font.bold = True
    subtitle.text_frame.paragraphs[0].font.size = Pt(20)

def add_content_slide(prs, title_text, bullet_points):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = title_text
    
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    
    for i, pt in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
            p.text = pt
        else:
            p = tf.add_paragraph()
            p.text = pt
            p.level = 0
            
        p.font.size = Pt(20)

def main():
    prs = Presentation()
    
    # Slide 1: Title
    add_title_slide(prs, "Error Analysis of POS Taggers", "NLP Final Project\nImplementation, Survey, Analysis, and Innovations")
    
    # Slide 2: Survey Done (Literature Review)
    survey_points = [
        "Conducted comprehensive literature review of POS tagging architectures.",
        "NLTK (Averaged Perceptron): Statistical baseline, relies on Penn Treebank tagset mapping.",
        "spaCy (en_core_web_sm): CNN-based neural approach, robust and fast.",
        "Stanza: BiLSTM + CRF architecture trained directly on Universal Dependencies.",
        "BERT (HuggingFace): Contextual Transformer fine-tuned on UPOS.",
        "Key insight: Transformer architectures handle context well but struggle with OOV domain shifts."
    ]
    add_content_slide(prs, "Survey Done (Literature Review)", survey_points)
    
    # Slide 3: Implementation Part
    impl_points = [
        "Dataset: Universal Dependencies English EWT (~204,000 train tokens).",
        "Pipeline Architecture: Created a unified interface to evaluate 4 distinct tagging models.",
        "Integration of multiple taggers (NLTK, spaCy, Stanza, BERT) into a single rigorous test bench.",
        "Developed custom Data Loaders for CoNLL-U format datasets.",
        "Automated the metric tracking: precision, recall, and F1 across all Universal POS categories."
    ]
    add_content_slide(prs, "Implementation Part", impl_points)
    
    # Slide 4: Analysis
    analysis_points = [
        "Six Dimensions of Analysis:",
        "1. Overall Accuracy: BERT & Stanza outperformed statistical baselines.",
        "2. Per-Tag F1: Identified inherently hard tags (e.g., ADJ vs VERB vs NOUN).",
        "3. Confusion Matrices: Quantified which UPOS categories are frequently confused.",
        "4. Out of Vocabulary (OOV): Analyzed error rates on unseen words.",
        "5. Frequency & Position: Evaluated errors by word frequency and sentence position.",
        "6. Lexical Ambiguity: Assessed models' ability to resolve multi-tag words."
    ]
    add_content_slide(prs, "Detailed Analysis & Results", analysis_points)
    
    # Slide 5: Innovation Part (CRITICAL)
    innovation_points = [
        "1. Majority-Vote Ensemble: Improved overall accuracy by combining predictions from all four distinct architectures.",
        "2. Comprehensive Error Taxonomy: Manually categorized errors into: OOV, Lexical Ambiguity, Proper Nouns, and Punctuation.",
        "3. Hard-Case Extraction: Isolated challenging tokens where ALL four taggers failed simultaneously.",
        "4. Error Breakdown: Detailed cross-model overlap analysis (how many models failed per token).",
        "5. Dashboard App: A full-stack Next.js dashboard providing dynamic visual analytics over the results."
    ]
    add_content_slide(prs, "Project Innovations", innovation_points)
    
    # Slide 6: Conclusion
    conclusion = [
        "Successfully developed an end-to-end evaluation pipeline for POS Taggers.",
        "Revealed specific weaknesses in modern NLP models (e.g., lexical ambiguity constraints).",
        "The ensemble approach highlighted the value of model diversity.",
        "Provides a reusable, visually informative framework for linguistic error analysis."
    ]
    add_content_slide(prs, "Conclusion", conclusion)
    
    prs.save("POS_Tagger_Analysis_Presentation.pptx")
    print("Presentation generated successfully: POS_Tagger_Analysis_Presentation.pptx")

if __name__ == "__main__":
    main()
